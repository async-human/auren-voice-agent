from __future__ import annotations

import base64
import io
import json
from collections.abc import Callable

from docx import Document
from httpx import AsyncClient
from openpyxl import load_workbook
from pptx import Presentation

from app.services.artifact_generators import (
    render_document,
    render_presentation,
    render_spreadsheet,
)
from app.services.action_payloads import redact_arguments


def _internal_user_id(participant_token: str) -> str:
    payload = participant_token.split(".")[1]
    claims = json.loads(base64.urlsafe_b64decode(payload + "=" * (-len(payload) % 4)))
    return json.loads(claims["metadata"])["userId"]


def test_document_generators_create_valid_formats() -> None:
    content = "## Findings\n\nA concise paragraph.\n\n- First point\n- Second point"
    assert render_document("Research", content, "pdf").startswith(b"%PDF")
    assert b"<h1>Research</h1>" in render_document("Research", content, "html")
    assert render_document("Research", content, "md").startswith(b"# Research")
    docx = Document(io.BytesIO(render_document("Research", content, "docx")))
    assert docx.core_properties.title == "Research"
    assert any(paragraph.text == "First point" for paragraph in docx.paragraphs)


def test_spreadsheet_generator_is_valid_and_neutralizes_formula_injection() -> None:
    generated = render_spreadsheet(
        "Results",
        ["Name", "Value"],
        [["safe", 4], ["=HYPERLINK(\"https://evil.test\")", 9]],
        "xlsx",
    )
    workbook = load_workbook(io.BytesIO(generated), data_only=False)
    sheet = workbook.active
    assert sheet.freeze_panes == "A2"
    assert sheet["A3"].value.startswith("'=")
    assert sheet["A3"].data_type != "f"


def test_presentation_generator_creates_title_and_content_slides() -> None:
    generated = render_presentation(
        "Quarterly review",
        [{"title": "Highlights", "bullets": ["Growth", "Reliability"], "notes": "Discuss."}],
    )
    presentation = Presentation(io.BytesIO(generated))
    assert len(presentation.slides) == 2
    assert presentation.slides[0].shapes.title.text == "Quarterly review"
    assert presentation.slides[1].shapes.title.text == "Highlights"


def test_artifact_content_is_omitted_from_audit_arguments() -> None:
    document = redact_arguments(
        "create_document",
        {"title": "Private", "content": "SECRET REPORT", "format": "docx"},
    )
    spreadsheet = redact_arguments(
        "create_spreadsheet",
        {"title": "Private", "rows": [["SECRET CELL"]], "format": "xlsx"},
    )
    presentation = redact_arguments(
        "create_presentation",
        {"title": "Private", "slides": [{"bullets": ["SECRET SLIDE"]}]},
    )
    assert document is not None and "SECRET REPORT" not in json.dumps(document)
    assert spreadsheet is not None and "SECRET CELL" not in json.dumps(spreadsheet)
    assert presentation is not None and "SECRET SLIDE" not in json.dumps(presentation)


async def test_artifact_tool_list_and_authenticated_download_end_to_end(
    client: AsyncClient,
    auth_headers: dict[str, str],
) -> None:
    token = await client.post("/v1/voice/token", headers=auth_headers)
    user_id = _internal_user_id(token.json()["participantToken"])
    created = await client.post(
        "/v1/tools/invoke",
        json={
            "tool": "create_document",
            "user_id": user_id,
            "arguments": {
                "title": "Auren test report",
                "content": "## Summary\n\nVerified end to end.",
                "format": "docx",
            },
        },
    )
    assert created.status_code == 200
    assert created.json()["ok"] is True
    artifact = created.json()["data"]["artifact"]

    listed = await client.get("/v1/artifacts", headers=auth_headers)
    assert listed.status_code == 200
    assert [item["id"] for item in listed.json()["artifacts"]] == [artifact["id"]]

    downloaded = await client.get(artifact["download_url"], headers=auth_headers)
    assert downloaded.status_code == 200
    assert downloaded.content.startswith(b"PK")
    assert downloaded.headers["x-artifact-sha256"] == artifact["sha256"]
    assert "attachment" in downloaded.headers["content-disposition"]


async def test_artifact_download_is_scoped_to_its_owner(
    client: AsyncClient,
    auth_headers: dict[str, str],
    make_token: Callable[..., str],
) -> None:
    token = await client.post("/v1/voice/token", headers=auth_headers)
    user_id = _internal_user_id(token.json()["participantToken"])
    created = await client.post(
        "/v1/tools/invoke",
        json={
            "tool": "create_document",
            "user_id": user_id,
            "arguments": {"title": "Private", "content": "Owner only", "format": "txt"},
        },
    )
    artifact_id = created.json()["data"]["artifact"]["id"]
    other_headers = {"Authorization": f"Bearer {make_token(subject='user_clerk_bob')}"}
    response = await client.get(
        f"/v1/artifacts/{artifact_id}/download",
        headers=other_headers,
    )
    assert response.status_code == 404
