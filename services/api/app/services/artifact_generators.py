from __future__ import annotations

import csv
import html
import io
import re
from collections.abc import Sequence
from dataclasses import dataclass
from typing import Any

from docx import Document
from docx.enum.text import WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.opc.constants import RELATIONSHIP_TYPE
from docx.shared import Inches, Pt
from markdown_it import MarkdownIt
from markdown_it.token import Token
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    HRFlowable,
    Paragraph,
    Preformatted,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

SpreadsheetValue = str | int | float | bool | None


@dataclass(frozen=True)
class InlineSpan:
    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False
    href: str | None = None


@dataclass(frozen=True)
class ContentBlock:
    kind: str
    inlines: tuple[InlineSpan, ...] = ()
    level: int = 0
    ordered: bool = False
    index: int | None = None
    rows: tuple[tuple[tuple[InlineSpan, ...], ...], ...] = ()


_MARKDOWN = MarkdownIt(
    "commonmark",
    {"html": False, "linkify": False, "typographer": False},
).enable("table")


def render_document(title: str, content: str, output_format: str) -> bytes:
    if output_format == "txt":
        return f"{title}\n{'=' * len(title)}\n\n{content.strip()}\n".encode()
    if output_format == "md":
        return f"# {title}\n\n{content.strip()}\n".encode()
    if output_format == "html":
        return _render_html(title, content).encode()
    if output_format == "docx":
        return _render_docx(title, content)
    if output_format == "pdf":
        return _render_pdf(title, content)
    raise ValueError(f"Unsupported document format: {output_format}")


def render_spreadsheet(
    title: str,
    columns: list[str],
    rows: list[list[SpreadsheetValue]],
    output_format: str,
) -> bytes:
    safe_rows = [[_safe_spreadsheet_value(value) for value in row] for row in rows]
    if output_format == "csv":
        stream = io.StringIO(newline="")
        writer = csv.writer(stream)
        writer.writerow(columns)
        writer.writerows(safe_rows)
        return stream.getvalue().encode("utf-8-sig")
    if output_format != "xlsx":
        raise ValueError(f"Unsupported spreadsheet format: {output_format}")

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = _worksheet_title(title)
    sheet.freeze_panes = "A2"
    sheet.append(columns)
    for row in safe_rows:
        sheet.append(row)
    header_fill = PatternFill("solid", fgColor="153B38")
    for cell in sheet[1]:
        cell.font = Font(color="FFFFFF", bold=True)
        cell.fill = header_fill
        cell.alignment = Alignment(vertical="center")
    sheet.auto_filter.ref = sheet.dimensions
    for index, column in enumerate(columns, start=1):
        values = [str(column), *(str(row[index - 1]) for row in safe_rows if index <= len(row))]
        sheet.column_dimensions[get_column_letter(index)].width = min(
            48,
            max(10, max((len(value) for value in values), default=10) + 2),
        )
    stream = io.BytesIO()
    workbook.save(stream)
    return stream.getvalue()


def render_presentation(title: str, slides: Sequence[dict[str, Any]]) -> bytes:
    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333)
    presentation.slide_height = PptxInches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Prepared by Auren"

    for item in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(item.get("title") or "Section")
        frame = slide.placeholders[1].text_frame
        frame.clear()
        bullets = item.get("bullets") or []
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.size = PptxPt(22)
        notes = item.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _render_docx(title: str, content: str) -> bytes:
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)
    document.core_properties.title = title
    document.add_heading(title, level=0)
    for block in _markdown_blocks(content):
        if block.kind == "table":
            _add_docx_table(document, block)
            continue
        if block.kind == "rule":
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.space_after = Pt(6)
            border = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:color"), "D7DFDD")
            border.append(bottom)
            paragraph._p.get_or_add_pPr().append(border)
            continue
        if block.kind == "heading":
            paragraph = document.add_heading(level=min(block.level, 3))
        elif block.kind == "list_item":
            paragraph = document.add_paragraph(
                style="List Number" if block.ordered else "List Bullet"
            )
            paragraph.paragraph_format.left_indent = Inches(
                max(block.level - 1, 0) * 0.22
            )
        elif block.kind == "quote":
            paragraph = document.add_paragraph(style="Intense Quote")
        elif block.kind == "code":
            paragraph = document.add_paragraph(style="No Spacing")
            shading = OxmlElement("w:shd")
            shading.set(qn("w:fill"), "F2F5F4")
            paragraph._p.get_or_add_pPr().append(shading)
        else:
            paragraph = document.add_paragraph()
        _add_docx_inlines(paragraph, block.inlines)
    styles = document.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5)
    stream = io.BytesIO()
    document.save(stream)
    return stream.getvalue()


def _render_pdf(title: str, content: str) -> bytes:
    stream = io.BytesIO()
    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "AurenBody",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10.2,
        leading=15,
        alignment=TA_LEFT,
        spaceAfter=7,
    )
    heading = ParagraphStyle(
        "AurenHeading",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=14,
        leading=18,
        textColor="#153B38",
        spaceBefore=10,
        spaceAfter=6,
    )
    quote = ParagraphStyle(
        "AurenQuote",
        parent=body,
        leftIndent=10,
        borderColor="#8BE9D5",
        borderWidth=1,
        borderPadding=7,
        backColor="#F3F8F6",
        textColor="#38504B",
    )
    code = ParagraphStyle(
        "AurenCode",
        parent=body,
        fontName="Courier",
        fontSize=8.8,
        leading=12,
        leftIndent=7,
        rightIndent=7,
        borderPadding=7,
        backColor="#F2F5F4",
    )
    story: list[Any] = [Paragraph(html.escape(title), styles["Title"]), Spacer(1, 5 * mm)]
    for block in _markdown_blocks(content):
        if block.kind == "heading":
            block_heading = ParagraphStyle(
                f"AurenHeading{block.level}",
                parent=heading,
                fontSize=max(11.5, 16 - block.level),
                leading=max(15, 20 - block.level),
            )
            story.append(Paragraph(_pdf_inline_markup(block.inlines), block_heading))
        elif block.kind == "list_item":
            list_style = ParagraphStyle(
                f"AurenList{block.level}",
                parent=body,
                leftIndent=10 + max(block.level - 1, 0) * 13,
                firstLineIndent=0,
                bulletIndent=max(block.level - 1, 0) * 13,
                spaceAfter=4,
            )
            marker = f"{block.index}." if block.ordered else "•"
            story.append(
                Paragraph(_pdf_inline_markup(block.inlines), list_style, bulletText=marker)
            )
        elif block.kind == "quote":
            story.append(Paragraph(_pdf_inline_markup(block.inlines), quote))
        elif block.kind == "code":
            story.append(Preformatted(_plain_text(block.inlines), code))
        elif block.kind == "rule":
            story.append(
                HRFlowable(
                    width="100%",
                    thickness=0.6,
                    color="#D7DFDD",
                    spaceBefore=5,
                    spaceAfter=8,
                )
            )
        elif block.kind == "table":
            story.append(_pdf_table(block, body))
        else:
            story.append(Paragraph(_pdf_inline_markup(block.inlines), body))
    document = SimpleDocTemplate(
        stream,
        pagesize=A4,
        rightMargin=18 * mm,
        leftMargin=18 * mm,
        topMargin=17 * mm,
        bottomMargin=17 * mm,
        title=title,
    )
    document.build(story)
    return stream.getvalue()


def _render_html(title: str, content: str) -> str:
    rendered = _MARKDOWN.render(content.strip())
    return "".join(
        [
            "<!doctype html><html><head><meta charset=\"utf-8\">",
            f"<title>{html.escape(title)}</title>",
            "<style>body{max-width:820px;margin:48px auto;padding:0 24px;"
            "font:16px/1.65 system-ui;color:#18211f}h1,h2,h3{color:#153b38}"
            "li{margin:.35em 0}code,pre{background:#f2f5f4}code{padding:.12em "
            ".3em;border-radius:4px}pre{padding:16px;overflow:auto;border-radius:10px}"
            "blockquote{margin-left:0;padding-left:18px;border-left:3px solid "
            "#8be9d5;color:#38504b}table{width:100%;border-collapse:collapse}"
            "th,td{padding:8px 10px;border:1px solid #d7dfdd;text-align:left}"
            "</style>",
            "</head><body>",
            f"<h1>{html.escape(title)}</h1>",
            rendered,
            "</body></html>",
        ]
    )


def _markdown_blocks(content: str) -> list[ContentBlock]:
    tokens = _MARKDOWN.parse(content)
    blocks: list[ContentBlock] = []
    list_stack: list[dict[str, int | bool]] = []
    item_stack: list[tuple[bool, int, int]] = []
    quote_depth = 0
    index = 0

    while index < len(tokens):
        token = tokens[index]
        if token.type in {"bullet_list_open", "ordered_list_open"}:
            start = int(token.attrGet("start") or 1)
            list_stack.append({"ordered": token.type == "ordered_list_open", "next": start})
        elif token.type in {"bullet_list_close", "ordered_list_close"}:
            list_stack.pop()
        elif token.type == "list_item_open" and list_stack:
            context = list_stack[-1]
            item_index = int(context["next"])
            context["next"] = item_index + 1
            item_stack.append((bool(context["ordered"]), item_index, len(list_stack)))
        elif token.type == "list_item_close" and item_stack:
            item_stack.pop()
        elif token.type == "blockquote_open":
            quote_depth += 1
        elif token.type == "blockquote_close":
            quote_depth = max(0, quote_depth - 1)
        elif token.type == "heading_open" and index + 1 < len(tokens):
            level = int(token.tag[1:]) if token.tag.startswith("h") else 2
            blocks.append(ContentBlock("heading", _inline_spans(tokens[index + 1]), level=level))
            index += 2
        elif token.type == "paragraph_open" and index + 1 < len(tokens):
            kind = "list_item" if item_stack else "quote" if quote_depth else "paragraph"
            ordered, item_index, level = item_stack[-1] if item_stack else (False, 0, 0)
            blocks.append(
                ContentBlock(
                    kind,
                    _inline_spans(tokens[index + 1]),
                    level=level,
                    ordered=ordered,
                    index=item_index if ordered else None,
                )
            )
            index += 2
        elif token.type in {"fence", "code_block"}:
            blocks.append(
                ContentBlock(
                    "code",
                    (InlineSpan(token.content.rstrip("\n"), code=True),),
                )
            )
        elif token.type == "hr":
            blocks.append(ContentBlock("rule"))
        elif token.type == "table_open":
            table, index = _parse_markdown_table(tokens, index)
            blocks.append(table)
        index += 1
    return blocks


def _inline_spans(token: Token) -> tuple[InlineSpan, ...]:
    children = token.children or []
    spans: list[InlineSpan] = []
    bold = italic = False
    href: str | None = None
    for child in children:
        if child.type == "strong_open":
            bold = True
        elif child.type == "strong_close":
            bold = False
        elif child.type == "em_open":
            italic = True
        elif child.type == "em_close":
            italic = False
        elif child.type == "link_open":
            href = child.attrGet("href")
        elif child.type == "link_close":
            href = None
        elif child.type in {"softbreak", "hardbreak"}:
            spans.append(InlineSpan("\n", bold=bold, italic=italic, href=href))
        elif child.type == "code_inline":
            spans.append(
                InlineSpan(
                    child.content,
                    bold=bold,
                    italic=italic,
                    code=True,
                    href=href,
                )
            )
        elif child.type == "image":
            spans.append(
                InlineSpan(
                    child.content or "Image",
                    bold=bold,
                    italic=italic,
                    href=child.attrGet("src"),
                )
            )
        elif child.type == "text":
            spans.append(InlineSpan(child.content, bold=bold, italic=italic, href=href))
    return tuple(spans)


def _parse_markdown_table(tokens: list[Token], start: int) -> tuple[ContentBlock, int]:
    rows: list[tuple[tuple[InlineSpan, ...], ...]] = []
    row: list[tuple[InlineSpan, ...]] | None = None
    cell: tuple[InlineSpan, ...] = ()
    index = start + 1
    while index < len(tokens) and tokens[index].type != "table_close":
        token = tokens[index]
        if token.type == "tr_open":
            row = []
        elif token.type in {"th_open", "td_open"}:
            cell = ()
        elif token.type == "inline":
            cell = _inline_spans(token)
        elif token.type in {"th_close", "td_close"} and row is not None:
            row.append(cell)
        elif token.type == "tr_close" and row is not None:
            rows.append(tuple(row))
            row = None
        index += 1
    return ContentBlock("table", rows=tuple(rows)), index


def _add_docx_inlines(paragraph: Any, spans: Sequence[InlineSpan]) -> None:
    for span in spans:
        pieces = span.text.split("\n")
        for piece_index, piece in enumerate(pieces):
            if piece:
                run = paragraph.add_run(piece)
                run.bold = span.bold
                run.italic = span.italic
                if span.code:
                    run.font.name = "Aptos Mono"
                if span.href:
                    run.font.underline = True
                    run.font.color.rgb = None
                    relationship = paragraph.part.relate_to(
                        span.href,
                        RELATIONSHIP_TYPE.HYPERLINK,
                        is_external=True,
                    )
                    hyperlink = OxmlElement("w:hyperlink")
                    hyperlink.set(qn("r:id"), relationship)
                    hyperlink.append(run._r)
                    paragraph._p.append(hyperlink)
            if piece_index < len(pieces) - 1:
                paragraph.add_run().add_break(WD_BREAK.LINE)


def _add_docx_table(document: Any, block: ContentBlock) -> None:
    if not block.rows:
        return
    column_count = max(len(row) for row in block.rows)
    table = document.add_table(rows=len(block.rows), cols=column_count)
    table.style = "Light Shading Accent 1"
    for row_index, row in enumerate(block.rows):
        for column_index, spans in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = ""
            paragraph = cell.paragraphs[0]
            if row_index == 0:
                spans = tuple(
                    InlineSpan(span.text, True, span.italic, span.code, span.href)
                    for span in spans
                )
            _add_docx_inlines(paragraph, spans)


def _pdf_inline_markup(spans: Sequence[InlineSpan]) -> str:
    parts: list[str] = []
    for span in spans:
        value = html.escape(span.text).replace("\n", "<br/>")
        if span.code:
            value = f'<font name="Courier">{value}</font>'
        if span.italic:
            value = f"<i>{value}</i>"
        if span.bold:
            value = f"<b>{value}</b>"
        if span.href:
            value = f'<a href="{html.escape(span.href, quote=True)}" color="#216E64">{value}</a>'
        parts.append(value)
    return "".join(parts)


def _plain_text(spans: Sequence[InlineSpan]) -> str:
    return "".join(span.text for span in spans)


def _pdf_table(block: ContentBlock, body: ParagraphStyle) -> Table:
    column_count = max((len(row) for row in block.rows), default=1)
    data: list[list[Paragraph]] = []
    for row in block.rows:
        cells = [Paragraph(_pdf_inline_markup(cell), body) for cell in row]
        cells.extend(Paragraph("", body) for _ in range(column_count - len(cells)))
        data.append(cells)
    table = Table(
        data,
        colWidths=[174 * mm / column_count] * column_count,
        repeatRows=1,
    )
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), "#E7F2EF"),
                ("TEXTCOLOR", (0, 0), (-1, 0), "#153B38"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.45, "#D7DFDD"),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )
    return table


def _safe_spreadsheet_value(value: SpreadsheetValue) -> SpreadsheetValue:
    if isinstance(value, str) and value.lstrip().startswith(("=", "+", "-", "@")):
        return "'" + value
    return value


def _worksheet_title(title: str) -> str:
    cleaned = re.sub(r"[\\/*?:\[\]]", " ", title).strip() or "Report"
    return cleaned[:31]
